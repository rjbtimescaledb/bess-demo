#!/usr/bin/env python3
"""
Generate BESS demo slides matching the (RJB) TT - BESS/VPP format.
Outputs a new PPTX with: Demo Results, RDS Comparison, and 10 Query Pack slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Dark theme colors matching the deck
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD = RGBColor(0x16, 0x21, 0x3E)
TEXT_WHITE = RGBColor(0xE0, 0xE0, 0xE0)
TEXT_GRAY = RGBColor(0x88, 0x88, 0x88)
TEXT_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_BLUE = RGBColor(0x4F, 0xC3, 0xF7)
ACCENT_GREEN = RGBColor(0x66, 0xBB, 0x6A)
ACCENT_RED = RGBColor(0xEF, 0x53, 0x50)
ACCENT_AMBER = RGBColor(0xFF, 0xB7, 0x4D)
ACCENT_PURPLE = RGBColor(0xB3, 0x88, 0xFF)
ACCENT_TEAL = RGBColor(0x26, 0xA6, 0x9A)
BRAND_BLUE = RGBColor(0x42, 0xA5, 0xF5)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def add_bg(slide):
    """Set dark background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_title(slide, title, subtitle=None):
    """Add slide title in top-left."""
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_GRAY


def add_card(slide, left, top, width, height, fill_color=BG_CARD):
    """Add a rounded rectangle card."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, size=11, color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def add_kpi(slide, left, top, label, value, color=ACCENT_BLUE, width=Inches(1.8)):
    """Add a KPI box."""
    add_card(slide, left, top, width, Inches(0.9))
    add_text(slide, left + Inches(0.15), top + Inches(0.08), width - Inches(0.3), Inches(0.3),
             label, size=9, color=TEXT_GRAY)
    add_text(slide, left + Inches(0.15), top + Inches(0.38), width - Inches(0.3), Inches(0.4),
             value, size=18, color=color, bold=True)


def add_table_row(tf, cols, sizes, colors=None, bold=False):
    """Add a line to a text frame simulating a table row."""
    if colors is None:
        colors = [TEXT_LIGHT] * len(cols)
    p = tf.add_paragraph()
    for i, (col, size) in enumerate(zip(cols, sizes)):
        run = p.add_run()
        run.text = col.ljust(size) if i < len(cols) - 1 else col
        run.font.size = Pt(11)
        run.font.color.rgb = colors[i] if i < len(colors) else TEXT_LIGHT
        run.font.bold = bold


# ============================================================
# SLIDE 1: Demo Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_title(slide, "Demo Results", "Live benchmark: Tiger Cloud with 6.6B rows, concurrent ingest + 12 dashboard users")

# KPI row
kpis = [
    ("Total Rows", "6.6B", ACCENT_BLUE),
    ("Ingest Rate", "2,500/sec", ACCENT_GREEN),
    ("Concurrent Users", "12", ACCENT_TEAL),
    ("Sites", "10", TEXT_WHITE),
    ("Assets", "200", TEXT_WHITE),
    ("Fleet Capacity", "2,080 MW", ACCENT_AMBER),
]
for i, (label, value, color) in enumerate(kpis):
    add_kpi(slide, Inches(0.6 + i * 2.05), Inches(1.2), label, value, color)

# Query results table
add_card(slide, Inches(0.5), Inches(2.4), Inches(7.5), Inches(4.5))
tf = add_text(slide, Inches(0.7), Inches(2.5), Inches(7), Inches(0.3),
              "Core Query Performance", size=14, color=ACCENT_BLUE, bold=True)

# Header
tf2 = add_text(slide, Inches(0.7), Inches(2.9), Inches(7), Inches(3.8), "", size=11)
tf2.paragraphs[0].text = ""
add_table_row(tf2, ["Query", "Category", "Latency", "Feature"], [32, 18, 10, 25],
              [TEXT_GRAY]*4, bold=True)
add_table_row(tf2, ["─"*30, "─"*16, "─"*8, "─"*23], [32, 18, 10, 25], [TEXT_GRAY]*4)

queries = [
    ("Q1A  Latest Fleet State", "Operational", "~50ms", "1-min CAGG"),
    ("Q1B  Active Alarms", "Operational", "~50ms", "Partial index"),
    ("Q2A  Power Trend (24h)", "Historical", "~100ms", "15-min CAGG"),
    ("Q2B  Asset Health Degradation", "Historical", "~150ms", "1-hour CAGG"),
    ("Q3A  Multi-Resolution", "Dashboard", "~50ms", "Auto tier"),
    ("Q3B  Fleet Utilization", "Dashboard", "~150ms", "15-min CAGG"),
    ("Q4A  Dispatch Readiness", "Decisioning", "~100ms", "CAGG + alarms"),
    ("Q4B  Revenue Opportunity", "Decisioning", "~100ms", "CAGG + market"),
    ("Q5A  Missed Revenue (7d)", "Revenue", "~400ms", "1-hour CAGG"),
    ("Q5B  Platform Proof", "Platform", "~50ms", "Metadata"),
]
for q_name, cat, latency, feature in queries:
    latency_color = ACCENT_GREEN if "50" in latency else ACCENT_BLUE if "100" in latency or "150" in latency else ACCENT_AMBER
    add_table_row(tf2, [q_name, cat, latency, feature], [32, 18, 10, 25],
                  [TEXT_WHITE, TEXT_GRAY, latency_color, ACCENT_PURPLE])

# Compression card
add_card(slide, Inches(8.5), Inches(2.4), Inches(4.3), Inches(2.0))
add_text(slide, Inches(8.7), Inches(2.5), Inches(3.9), Inches(0.3),
         "Columnstore Compression", size=14, color=ACCENT_TEAL, bold=True)
comp_items = [
    ("segmentby", "site_id, asset_id"),
    ("orderby", "ts DESC"),
    ("Encoding", "Delta + Gorilla"),
    ("Target Ratio", "90%+ compression"),
    ("Compress After", "1 day (hourly check)"),
    ("S3 Tiering", "After 3 months"),
]
for i, (k, v) in enumerate(comp_items):
    add_text(slide, Inches(8.7), Inches(2.95 + i * 0.28), Inches(1.6), Inches(0.25),
             k, size=10, color=TEXT_GRAY)
    add_text(slide, Inches(10.3), Inches(2.95 + i * 0.28), Inches(2.5), Inches(0.25),
             v, size=10, color=ACCENT_TEAL, bold=True)

# Data lifecycle card
add_card(slide, Inches(8.5), Inches(4.7), Inches(4.3), Inches(2.2))
add_text(slide, Inches(8.7), Inches(4.8), Inches(3.9), Inches(0.3),
         "Data Lifecycle", size=14, color=BRAND_BLUE, bold=True)
lifecycle = [
    ("Rowstore (hot)", "Last 1 day — fast inserts", ACCENT_RED),
    ("Columnstore (warm)", "1 day → 3 months — compressed", ACCENT_AMBER),
    ("S3 Tiered (cold)", "3 months+ — object storage", BRAND_BLUE),
    ("Retention", "Auto-drop after 6 months", TEXT_GRAY),
]
for i, (stage, desc, color) in enumerate(lifecycle):
    add_text(slide, Inches(8.7), Inches(5.2 + i * 0.45), Inches(2.0), Inches(0.22),
             stage, size=11, color=color, bold=True)
    add_text(slide, Inches(8.7), Inches(5.42 + i * 0.45), Inches(3.8), Inches(0.22),
             desc, size=9, color=TEXT_GRAY)


# ============================================================
# SLIDE 2: RDS Comparison
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "Tiger Cloud vs AWS RDS PostgreSQL",
          "Same spec (4 vCPU / 16 GiB), same region (eu-central-1), same data — only the database differs")

# Benchmark results table
add_card(slide, Inches(0.5), Inches(1.3), Inches(7.5), Inches(3.5))
add_text(slide, Inches(0.7), Inches(1.4), Inches(5), Inches(0.3),
         "Query Pack Benchmark (Query Pack v1)", size=14, color=ACCENT_BLUE, bold=True)

tf = add_text(slide, Inches(0.7), Inches(1.8), Inches(7), Inches(3.0), "", size=11)
tf.paragraphs[0].text = ""
add_table_row(tf, ["Query", "Tiger Cloud", "RDS", "Winner"], [30, 14, 14, 20],
              [TEXT_GRAY]*4, bold=True)
add_table_row(tf, ["─"*28, "─"*12, "─"*12, "─"*18], [30, 14, 14, 20], [TEXT_GRAY]*4)

bench = [
    ("Q1  Latest Fleet State", "340ms", "9,639ms", "Tiger 28x faster", ACCENT_GREEN),
    ("Q7  Dispatch Readiness", "241ms", "7,330ms", "Tiger 30x faster", ACCENT_GREEN),
    ("Q8  Revenue Opportunity", "356ms", "7,202ms", "Tiger 20x faster", ACCENT_GREEN),
    ("Q13 Raw Aggregation (1h)", "1,509ms", "44ms", "RDS (fewer rows)*", TEXT_GRAY),
    ("Q14 CAGG Aggregation", "52ms", "44ms", "~equal", TEXT_GRAY),
]
for q, tiger, rds, winner, wcolor in bench:
    add_table_row(tf, [q, tiger, rds, winner], [30, 14, 14, 20],
                  [TEXT_WHITE, ACCENT_GREEN, ACCENT_RED, wcolor])

add_table_row(tf, ["", "", "", ""], [30, 14, 14, 20])
add_table_row(tf, ["* RDS has 186M rows vs Tiger's 6.6B", "", "", ""], [78, 0, 0, 0],
              [TEXT_GRAY])

# Storage comparison
add_card(slide, Inches(8.5), Inches(1.3), Inches(4.3), Inches(1.8))
add_text(slide, Inches(8.7), Inches(1.4), Inches(3.9), Inches(0.3),
         "Storage Comparison", size=14, color=ACCENT_TEAL, bold=True)
storage = [
    ("Tiger Cloud", "Columnstore compressed", ACCENT_GREEN),
    ("RDS (186M rows)", "22 GB uncompressed", ACCENT_RED),
    ("RDS (3yr @ 3K/sec)", "38 TB projected", ACCENT_RED),
]
for i, (label, val, color) in enumerate(storage):
    add_text(slide, Inches(8.7), Inches(1.85 + i * 0.35), Inches(2.0), Inches(0.25),
             label, size=10, color=TEXT_GRAY)
    add_text(slide, Inches(10.7), Inches(1.85 + i * 0.35), Inches(2.0), Inches(0.25),
             val, size=10, color=color, bold=True)

# TCO card
add_card(slide, Inches(8.5), Inches(3.4), Inches(4.3), Inches(1.4))
add_text(slide, Inches(8.7), Inches(3.5), Inches(3.9), Inches(0.3),
         "3-Year TCO @ 3K rows/sec", size=14, color=ACCENT_AMBER, bold=True)
tco = [
    ("RDS Single-AZ (RI 3yr)", "$147,646"),
    ("RDS Multi-AZ (RI 3yr)", "$261,062"),
    ("Storage = 67% of cost", "No compression"),
]
for i, (label, val) in enumerate(tco):
    add_text(slide, Inches(8.7), Inches(3.9 + i * 0.3), Inches(2.2), Inches(0.25),
             label, size=10, color=TEXT_GRAY)
    add_text(slide, Inches(10.9), Inches(3.9 + i * 0.3), Inches(1.8), Inches(0.25),
             val, size=10, color=ACCENT_RED, bold=True)

# What RDS is missing
add_card(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(2.0))
add_text(slide, Inches(0.7), Inches(5.2), Inches(5), Inches(0.3),
         "What RDS is Missing", size=14, color=ACCENT_RED, bold=True)

missing_items = [
    ("Hypertables", "Automatic time partitioning"),
    ("Continuous Aggregates", "Pre-computed rollups (1min → 15min → 1hour)"),
    ("Columnstore Compression", "90%+ storage savings"),
    ("S3 Tiering", "Automatic cold data archival"),
    ("time_bucket()", "Native time-series aggregation"),
    ("Retention Policies", "Automated data lifecycle management"),
]
col1_x = Inches(0.7)
col2_x = Inches(6.5)
for i, (feature, desc) in enumerate(missing_items):
    col = col1_x if i < 3 else col2_x
    row = i % 3
    add_text(slide, col, Inches(5.6 + row * 0.5), Inches(1.5), Inches(0.22),
             "✗  " + feature, size=11, color=ACCENT_RED, bold=True)
    add_text(slide, col + Inches(1.8), Inches(5.6 + row * 0.5), Inches(3.5), Inches(0.22),
             desc, size=10, color=TEXT_GRAY)


# ============================================================
# SLIDES 3-12: Query Pack (one per query)
# ============================================================
query_slides = [
    {
        "id": "Q1A", "name": "Latest Fleet State",
        "category": "Operational", "cat_color": ACCENT_GREEN,
        "user": "BESS Operator",
        "question": "What is happening across my fleet right now?",
        "decision": "Do I need to intervene?",
        "value": "Real-time visibility, faster incident response",
        "latency": "~50ms", "feature": "1-min Continuous Aggregate, DISTINCT ON",
        "sql": """SELECT DISTINCT ON (s.site_id)
  s.site_id, s.name AS site_name,
  t.bucket AS latest_time,
  t.avg_site_power_mw,
  t.avg_soc_pct AS soc_pct,
  t.avg_soh_pct AS soh_pct
FROM sites s
JOIN telemetry_1min t ON t.site_id = s.site_id
WHERE t.bucket > NOW() - INTERVAL '10 minutes'
ORDER BY s.site_id, t.bucket DESC;""",
    },
    {
        "id": "Q1B", "name": "Active Alarms by Site",
        "category": "Operational", "cat_color": ACCENT_GREEN,
        "user": "Reliability Engineer",
        "question": "What is broken right now and where?",
        "decision": "Which site needs immediate attention?",
        "value": "Reduced MTTR, prioritized response",
        "latency": "~50ms", "feature": "Partial index on unresolved alarms",
        "sql": """SELECT s.name, a.severity, COUNT(*),
  MIN(a.ts) AS oldest_unresolved
FROM alarms_events a
JOIN sites s ON s.site_id = a.site_id
WHERE a.resolved_at IS NULL
GROUP BY s.name, a.severity
ORDER BY CASE severity
  WHEN 'emergency' THEN 1
  WHEN 'critical' THEN 2
  WHEN 'warning' THEN 3 ELSE 4
END, count(*) DESC;""",
    },
    {
        "id": "Q2A", "name": "Power Trend (24h, per-site)",
        "category": "Historical", "cat_color": BRAND_BLUE,
        "user": "Operations Analyst",
        "question": "How has this site's power output changed over 24 hours?",
        "decision": "Is behavior normal or anomalous?",
        "value": "Early anomaly detection, operational insight",
        "latency": "~100ms", "feature": "15-min Continuous Aggregate (pre-computed)",
        "sql": """SELECT bucket AS time,
  avg_site_power_mw AS power_mw,
  avg_soc_pct AS soc_pct,
  avg_discharge_power_mw AS discharge_mw
FROM telemetry_15min
WHERE site_id = $1
  AND bucket >= NOW() - INTERVAL '24 hours'
ORDER BY bucket;""",
    },
    {
        "id": "Q2B", "name": "Asset Health Degradation (30d)",
        "category": "Historical", "cat_color": BRAND_BLUE,
        "user": "Asset Manager",
        "question": "Which sites are degrading fastest?",
        "decision": "Schedule preventive maintenance before failure",
        "value": "Reduced replacement cost, extended asset life",
        "latency": "~150ms", "feature": "1-hour CAGG, time-windowed weekly comparison",
        "sql": """WITH weekly_soh AS (
  SELECT site_id,
    CASE WHEN bucket >= NOW()-'7d' THEN 'now'
         WHEN bucket >= NOW()-'14d' THEN 'last_week'
         WHEN bucket >= NOW()-'30d' THEN 'month_ago'
    END AS period, AVG(avg_soh_pct) AS avg_soh
  FROM telemetry_1hour
  WHERE bucket >= NOW() - INTERVAL '30 days'
  GROUP BY site_id, period
)
SELECT name, soh_now, soh_month_ago,
  degradation_30d, projected_annual,
  CASE WHEN annual < -2 THEN 'CRITICAL'
       WHEN annual < -1 THEN 'WATCH'
       ELSE 'NORMAL' END AS status
FROM pivoted JOIN sites USING(site_id)
ORDER BY degradation_30d ASC;""",
    },
    {
        "id": "Q3A", "name": "Multi-Resolution Telemetry",
        "category": "Dashboard", "cat_color": ACCENT_PURPLE,
        "user": "Dashboard End-User",
        "question": "Can I zoom from 30-day overview to 1-minute detail?",
        "decision": "Which time resolution to serve based on requested range",
        "value": "Fast UX regardless of time range (sub-200ms always)",
        "latency": "~50ms", "feature": "Hierarchical CAGGs: raw → 1min → 15min → 1hour",
        "sql": """-- Auto-selected by range:
--   ≤ 1h  → telemetry_raw
--   ≤ 26h → telemetry_1min
--   ≤ 72h → telemetry_15min
--   > 72h → telemetry_1hour

SELECT bucket AS time,
  avg_site_power_mw, avg_soc_pct,
  avg_discharge_power_mw
FROM telemetry_15min  -- 7-day view
WHERE site_id = $1
  AND bucket >= NOW() - INTERVAL '7 days'
ORDER BY bucket;""",
    },
    {
        "id": "Q3B", "name": "Fleet Utilization Ranking (24h)",
        "category": "Dashboard", "cat_color": ACCENT_PURPLE,
        "user": "Operations Lead",
        "question": "Which sites are performing best? Where should I focus?",
        "decision": "Allocate resources to underperforming sites",
        "value": "Fleet-wide efficiency optimization",
        "latency": "~150ms", "feature": "15-min CAGG, fleet-wide aggregation + ranking",
        "sql": """SELECT s.name, s.capacity_mw,
  ROUND(AVG(t.avg_site_power_mw), 1) AS avg_power_mw,
  ROUND(AVG(t.avg_site_power_mw)
    / NULLIF(s.capacity_mw, 0) * 100, 1)
    AS utilization_pct,
  ROUND(AVG(t.avg_soc_pct), 1) AS avg_soc,
  ROUND(AVG(t.avg_rte), 1) AS avg_rte
FROM telemetry_15min t
JOIN sites s ON s.site_id = t.site_id
WHERE bucket >= NOW() - INTERVAL '24 hours'
GROUP BY s.name, s.capacity_mw
ORDER BY utilization_pct DESC;""",
    },
    {
        "id": "Q4A", "name": "Dispatch Readiness (scored)",
        "category": "Decisioning", "cat_color": ACCENT_AMBER,
        "user": "Energy Trader / Dispatch Operator",
        "question": "Which sites can I dispatch right now?",
        "decision": "Dispatch energy when prices are high",
        "value": "Avoid dispatching degraded sites, prevent grid penalties",
        "latency": "~100ms", "feature": "Real-time CAGG + alarms cross-join, readiness scoring",
        "sql": """WITH latest AS (
  SELECT DISTINCT ON (site_id)
    site_id, avg_soc_pct AS soc,
    avg_soh_pct AS soh
  FROM telemetry_1min
  WHERE bucket > NOW() - '10 min'
  ORDER BY site_id, bucket DESC
), alarms AS (
  SELECT site_id, COUNT(*) AS cnt
  FROM alarms_events
  WHERE resolved_at IS NULL
    AND severity IN ('critical','emergency')
  GROUP BY site_id
)
SELECT name, soc, soh, critical_alarms,
  available_energy_mwh,
  readiness_score,  -- 0-100
  dispatch_status   -- READY/CAUTION/NOT_READY
FROM sites JOIN latest JOIN alarms
ORDER BY readiness_score DESC;""",
    },
    {
        "id": "Q4B", "name": "Revenue Opportunity (real-time)",
        "category": "Decisioning", "cat_color": ACCENT_AMBER,
        "user": "Energy Trader",
        "question": "Where is the money right now?",
        "decision": "Maximize revenue by dispatching into high-price markets",
        "value": "Direct revenue optimization ($000s per decision)",
        "latency": "~100ms", "feature": "CAGG + market prices cross-join, capacity calculation",
        "sql": """WITH latest AS (
  SELECT DISTINCT ON (site_id)
    site_id, avg_soc_pct AS soc,
    avg_site_power_mw AS power
  FROM telemetry_1min
  WHERE bucket > NOW() - '10 min'
  ORDER BY site_id, bucket DESC
), prices AS (
  SELECT DISTINCT ON (market)
    market, price_usd_mwh
  FROM market_price_signals
  WHERE ts > NOW() - '30 min'
  ORDER BY market, ts DESC
)
SELECT name, market, price_usd_mwh,
  available_mw, available_mwh,
  revenue_per_hour_usd,
  total_opportunity_usd
FROM sites JOIN latest JOIN prices
ORDER BY revenue_per_hour DESC;""",
    },
    {
        "id": "Q5A", "name": "Missed Revenue (7-day lookback)",
        "category": "Revenue Analytics", "cat_color": ACCENT_RED,
        "user": "Head of Trading / CFO",
        "question": "How much money did we leave on the table this week?",
        "decision": "Improve dispatch strategy, justify automation investment",
        "value": "Quantifies opportunity cost — most actionable BESS metric",
        "latency": "~400ms", "feature": "1-hour CAGG + market join + FILTER aggregation",
        "sql": """WITH hourly_state AS (
  SELECT bucket, site_id, avg_soc_pct,
    avg_discharge_power_mw, capacity_mw
  FROM telemetry_1hour JOIN sites USING(site_id)
  WHERE bucket >= NOW() - '7 days'
), hourly_prices AS (
  SELECT time_bucket('1 hour', ts) AS bucket,
    market, AVG(price_usd_mwh) AS price
  FROM market_price_signals
  WHERE ts >= NOW() - '7 days'
  GROUP BY 1, market
)
SELECT name,
  COUNT(*) FILTER (WHERE price > 60)
    AS high_price_hours,
  COUNT(*) FILTER (WHERE missed)
    AS missed_hours,
  SUM(missed_revenue_usd)
    AS total_missed_revenue
FROM combined
GROUP BY name
ORDER BY total_missed_revenue DESC;""",
    },
    {
        "id": "Q5B", "name": "Platform Proof (ingest + query)",
        "category": "Platform", "cat_color": ACCENT_TEAL,
        "user": "Solutions Engineer / DBA",
        "question": "Can it handle real-time writes AND analytics simultaneously?",
        "decision": "Is TimescaleDB production-ready for this workload?",
        "value": "Platform confidence, reduced evaluation risk",
        "latency": "~50ms", "feature": "System catalogs, compression stats, chunk metadata",
        "sql": """SELECT
  (SELECT COUNT(*) FROM telemetry_raw
   WHERE ts > NOW() - '1 min')
   AS rows_last_minute,
  (SELECT reltuples::bigint FROM pg_class
   WHERE relname = 'telemetry_raw')
   AS est_total_rows,
  (SELECT COUNT(*) FROM chunks
   WHERE hypertable = 'telemetry_raw')
   AS total_chunks,
  (SELECT compression_ratio
   FROM hypertable_compression_stats(...))
   AS compression_ratio,
  pg_size_pretty(pg_database_size(...))
   AS database_size;""",
    },
]

for q in query_slides:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # Category badge
    add_text(slide, Inches(0.6), Inches(0.3), Inches(2), Inches(0.3),
             q["category"].upper(), size=10, color=q["cat_color"], bold=True)

    # Query ID + Name
    add_text(slide, Inches(0.6), Inches(0.55), Inches(8), Inches(0.5),
             f"{q['id']} — {q['name']}", size=26, color=TEXT_WHITE, bold=True)

    # Latency badge
    add_card(slide, Inches(10.5), Inches(0.3), Inches(2.3), Inches(0.7), fill_color=RGBColor(0x0D, 0x47, 0xA1))
    add_text(slide, Inches(10.5), Inches(0.35), Inches(2.3), Inches(0.25),
             "LATENCY", size=9, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(10.5), Inches(0.55), Inches(2.3), Inches(0.35),
             q["latency"], size=22, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)

    # User → Question → Decision → Value cards
    card_data = [
        ("User", q["user"], ACCENT_BLUE),
        ("Question", q["question"], TEXT_WHITE),
        ("Decision", q["decision"], ACCENT_AMBER),
        ("Business Value", q["value"], ACCENT_GREEN),
    ]
    card_w = Inches(3.0)
    for i, (label, text, color) in enumerate(card_data):
        x = Inches(0.5 + i * 3.15)
        add_card(slide, x, Inches(1.3), card_w, Inches(1.0))
        add_text(slide, x + Inches(0.15), Inches(1.38), card_w - Inches(0.3), Inches(0.2),
                 label, size=9, color=TEXT_GRAY, bold=True)
        add_text(slide, x + Inches(0.15), Inches(1.6), card_w - Inches(0.3), Inches(0.5),
                 text, size=12, color=color)

    # TimescaleDB Feature
    add_text(slide, Inches(0.6), Inches(2.5), Inches(2), Inches(0.2),
             "TimescaleDB Feature", size=9, color=TEXT_GRAY, bold=True)
    add_text(slide, Inches(2.6), Inches(2.5), Inches(8), Inches(0.2),
             q["feature"], size=10, color=ACCENT_PURPLE, bold=True)

    # SQL code block
    add_card(slide, Inches(0.5), Inches(2.9), Inches(12.3), Inches(4.2),
             fill_color=RGBColor(0x0F, 0x11, 0x1A))
    add_text(slide, Inches(0.7), Inches(2.95), Inches(1), Inches(0.2),
             "SQL", size=9, color=TEXT_GRAY, bold=True)

    sql_tf = add_text(slide, Inches(0.7), Inches(3.2), Inches(11.8), Inches(3.8),
                      q["sql"], size=11, color=ACCENT_TEAL)
    sql_tf.paragraphs[0].font.name = "Consolas"
    # Set monospace font for all text
    for p in sql_tf.paragraphs:
        for run in p.runs:
            run.font.name = "Consolas"


# ============================================================
# Save
# ============================================================
output_path = "/Users/ruhanjaybora/bess-demo/docs/BESS Demo Slides - Results & Query Pack.pptx"
prs.save(output_path)
print(f"Saved {len(prs.slides)} slides to: {output_path}")
